"""Build the Agent Austin presentation (dark theme, 16:9).

Run:
    uv run python docs/presentation/build_ppt.py
    # or
    python3 docs/presentation/build_ppt.py

Outputs docs/presentation/agent-austin.pptx.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).parent / "agent-austin.pptx"

BG = RGBColor(0x1A, 0x1A, 0x2E)
CARD = RGBColor(0x16, 0x21, 0x3E)
ACCENT_DEEP = RGBColor(0x0F, 0x34, 0x60)
CYAN = RGBColor(0x00, 0xD2, 0xFF)
PURPLE = RGBColor(0x7B, 0x2F, 0xF7)
CORAL = RGBColor(0xFF, 0x6B, 0x6B)
YELLOW = RGBColor(0xFF, 0xD9, 0x3D)
GREEN = RGBColor(0x6B, 0xCB, 0x77)
TEXT = RGBColor(0xE0, 0xE0, 0xE0)
MUTED = RGBColor(0x99, 0x99, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def paint_background(slide, color=BG):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    size=18,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT,
    font="Calibri",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_card(slide, left, top, width, height, *, color=CARD, radius=True):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(shape, left, top, width, height)
    card.line.fill.background()
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.shadow.inherit = False
    if radius and hasattr(card, "adjustments"):
        try:
            card.adjustments[0] = 0.08
        except Exception:
            pass
    return card


def add_accent_bar(slide, top, color=CYAN, left=Inches(0.6), width=Inches(0.12), height=Inches(0.5)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def add_footer(slide, page_num, total):
    add_text(
        slide,
        "Agent Austin",
        Inches(0.6),
        Inches(7.05),
        Inches(4),
        Inches(0.3),
        size=10,
        color=MUTED,
    )
    add_text(
        slide,
        f"{page_num} / {total}",
        Inches(12.0),
        Inches(7.05),
        Inches(0.8),
        Inches(0.3),
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


# ------- slide builders -------

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    paint_background(s)
    # Accent gradient band
    band = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(3.1), Inches(13.333), Inches(0.05)
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = CYAN
    add_text(
        s,
        "Agent Austin",
        Inches(0.6),
        Inches(2.3),
        Inches(12),
        Inches(1.2),
        size=64,
        bold=True,
        color=WHITE,
    )
    add_text(
        s,
        "An AI data-science agent for Austin 311 service requests",
        Inches(0.6),
        Inches(3.3),
        Inches(12),
        Inches(0.6),
        size=24,
        color=CYAN,
    )
    add_text(
        s,
        "Architecture · Agent Skills · Visualization · Use Cases",
        Inches(0.6),
        Inches(4.1),
        Inches(12),
        Inches(0.5),
        size=18,
        color=MUTED,
    )
    return s


def slide_agenda(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s, Inches(0.6), CYAN)
    add_text(
        s,
        "Agenda",
        Inches(0.85),
        Inches(0.55),
        Inches(10),
        Inches(0.7),
        size=36,
        bold=True,
        color=WHITE,
    )
    items = [
        ("1", "What is Agent Austin?", "The product and the data it sits on"),
        ("2", "Architecture", "How the pieces fit together"),
        ("3", "Agent Skills", "What the agent knows how to do"),
        ("4", "Visualization", "Plotly charts, dashboards, and reports"),
        ("5", "Use Cases", "Potholes, trash, code compliance, and more"),
        ("6", "Roadmap", "Where this is headed"),
    ]
    top = Inches(1.6)
    for i, (n, title, sub) in enumerate(items):
        y = top + Inches(i * 0.85)
        # number badge
        badge = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.85), y, Inches(0.55), Inches(0.55)
        )
        badge.line.fill.background()
        badge.fill.solid()
        badge.fill.fore_color.rgb = CYAN if i % 2 == 0 else PURPLE
        tf = badge.text_frame
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = n
        r.font.bold = True
        r.font.size = Pt(20)
        r.font.color.rgb = BG
        # title + subtitle
        add_text(s, title, Inches(1.65), y + Inches(0.02), Inches(10), Inches(0.45), size=22, bold=True, color=WHITE)
        add_text(s, sub, Inches(1.65), y + Inches(0.45), Inches(10), Inches(0.4), size=14, color=MUTED)
    return s


def slide_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s, Inches(0.6), CYAN)
    add_text(s, "What is Agent Austin?", Inches(0.85), Inches(0.55), Inches(12), Inches(0.7), size=36, bold=True, color=WHITE)
    add_text(
        s,
        "A chat-first data-science agent for the City of Austin's 311 service-request dataset.\nAsk in plain English — it downloads, analyzes, visualizes, and writes reports for you.",
        Inches(0.85),
        Inches(1.35),
        Inches(12),
        Inches(1.1),
        size=18,
        color=TEXT,
    )

    # Three stat cards
    stats = [
        ("~2.4M", "311 requests\n2014 → present", CYAN),
        ("Daily", "Delta-merged on\neach deploy", PURPLE),
        ("6", "Skills the agent\nknows how to run", CORAL),
    ]
    card_w = Inches(3.8)
    card_h = Inches(2.1)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_x = (Inches(13.333) - total_w) / 2
    for i, (val, label, color) in enumerate(stats):
        x = start_x + (card_w + gap) * i
        add_card(s, x, Inches(2.75), card_w, card_h)
        add_text(s, val, x, Inches(2.95), card_w, Inches(1.1), size=54, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, label, x, Inches(4.0), card_w, Inches(0.8), size=15, color=TEXT, align=PP_ALIGN.CENTER)

    # Data source line
    add_text(
        s,
        "Data: City of Austin Open Data · Socrata dataset xwdj-i9he · no API key required",
        Inches(0.85),
        Inches(5.3),
        Inches(12),
        Inches(0.5),
        size=14,
        color=MUTED,
    )
    return s


def _pill(slide, left, top, width, height, label, *, fill=CARD, text_color=WHITE, size=13, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.color.rgb = CYAN
    shape.line.width = Pt(0.75)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.shadow.inherit = False
    if hasattr(shape, "adjustments"):
        try:
            shape.adjustments[0] = 0.25
        except Exception:
            pass
    tf = shape.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    lines = label.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = text_color
    return shape


def _group_box(slide, left, top, width, height, title, *, color=CYAN):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.line.color.rgb = color
    box.line.width = Pt(1.25)
    box.line.dash_style = 7  # dash
    box.fill.background()
    box.shadow.inherit = False
    if hasattr(box, "adjustments"):
        try:
            box.adjustments[0] = 0.04
        except Exception:
            pass
    add_text(slide, title, left + Inches(0.2), top + Inches(0.08), width, Inches(0.35), size=12, bold=True, color=color)
    return box


def _connector(slide, x1, y1, x2, y2, *, color=CYAN, label=None):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # straight
    conn.line.color.rgb = color
    conn.line.width = Pt(1.5)
    if label:
        # midpoint label
        mx, my = (x1 + x2) / 2, (y1 + y2) / 2
        add_text(slide, label, mx - Inches(0.9), my - Inches(0.18), Inches(1.8), Inches(0.35), size=10, color=MUTED, align=PP_ALIGN.CENTER)
    return conn


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s, Inches(0.6), CYAN)
    add_text(s, "Architecture", Inches(0.85), Inches(0.55), Inches(12), Inches(0.7), size=36, bold=True, color=WHITE)
    add_text(s, "Browser → Next.js → FastAPI → Claude Agent SDK, backed by Postgres, a volume, and Socrata.", Inches(0.85), Inches(1.2), Inches(12), Inches(0.4), size=14, color=MUTED)

    # Three column groups
    # Column 1: Browser
    _group_box(s, Inches(0.6), Inches(1.9), Inches(2.6), Inches(4.6), "USER BROWSER", color=CORAL)
    _pill(s, Inches(0.85), Inches(3.6), Inches(2.1), Inches(1.0), "Next.js 16 UI\nAI Elements + shadcn", fill=CARD, size=12)

    # Column 2: Railway
    _group_box(s, Inches(3.5), Inches(1.9), Inches(6.3), Inches(4.6), "RAILWAY PLATFORM", color=CYAN)
    # Frontend service
    _group_box(s, Inches(3.7), Inches(2.4), Inches(2.9), Inches(1.3), "Frontend service", color=PURPLE)
    _pill(s, Inches(3.85), Inches(2.85), Inches(2.6), Inches(0.7), "Next.js App Router", fill=CARD, size=12)
    # Backend service
    _group_box(s, Inches(6.75), Inches(2.4), Inches(2.9), Inches(3.9), "Backend service", color=PURPLE)
    _pill(s, Inches(6.9), Inches(2.85), Inches(2.6), Inches(0.6), "FastAPI + uvicorn", fill=CARD, size=12)
    _pill(s, Inches(6.9), Inches(3.55), Inches(2.6), Inches(0.6), "Claude Agent SDK", fill=ACCENT_DEEP, size=12)
    _pill(s, Inches(6.9), Inches(4.25), Inches(2.6), Inches(0.6), "In-process MCP server", fill=ACCENT_DEEP, size=12)
    # Data at the bottom of the Railway column
    _pill(s, Inches(3.85), Inches(5.5), Inches(2.6), Inches(0.7), "PostgreSQL\nsessions / messages", fill=ACCENT_DEEP, size=11, bold=False)
    _pill(s, Inches(6.9), Inches(5.5), Inches(2.6), Inches(0.7), "Persistent Volume\n311 CSV · charts · reports", fill=ACCENT_DEEP, size=11, bold=False)

    # Column 3: External
    _group_box(s, Inches(10.1), Inches(1.9), Inches(2.6), Inches(4.6), "EXTERNAL", color=YELLOW)
    _pill(s, Inches(10.3), Inches(2.8), Inches(2.25), Inches(0.9), "Anthropic API\nClaude models", fill=CARD, size=12)
    _pill(s, Inches(10.3), Inches(4.1), Inches(2.25), Inches(0.9), "Socrata\nxwdj-i9he dataset", fill=CARD, size=12)

    # Connectors
    _connector(s, Inches(2.95), Inches(4.1), Inches(3.85), Inches(3.2), color=CYAN)  # browser → frontend
    _connector(s, Inches(6.45), Inches(3.15), Inches(6.9), Inches(3.15), color=CYAN)  # frontend → FastAPI (SSE)
    _connector(s, Inches(8.2), Inches(4.85), Inches(8.2), Inches(5.5), color=PURPLE)  # backend → volume
    _connector(s, Inches(5.15), Inches(3.55), Inches(5.15), Inches(5.5), color=PURPLE)  # frontend column → Postgres? actually backend → Postgres
    _connector(s, Inches(9.5), Inches(3.85), Inches(10.3), Inches(3.25), color=YELLOW)  # SDK → Anthropic
    _connector(s, Inches(9.5), Inches(4.55), Inches(10.3), Inches(4.55), color=YELLOW)  # backend → Socrata

    # Legend / notes
    add_text(
        s,
        "SSE streaming · Vercel AI SDK v6 protocol · JWT auth · delta-merge on startup",
        Inches(0.6),
        Inches(6.65),
        Inches(12),
        Inches(0.35),
        size=12,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    return s


def slide_skills_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s, Inches(0.6), PURPLE)
    add_text(s, "Agent Skills", Inches(0.85), Inches(0.55), Inches(12), Inches(0.7), size=36, bold=True, color=WHITE)
    add_text(
        s,
        "Markdown instructions under agent311/.claude/skills/ — the SDK picks them up automatically.",
        Inches(0.85),
        Inches(1.2),
        Inches(12),
        Inches(0.4),
        size=14,
        color=MUTED,
    )

    # 6 skill cards: 3 columns × 2 rows
    cards = [
        ("download-311-data", "Pulls Socrata CSV\npagination + delta merge", CYAN),
        ("analyze-311-data", "Exploratory stats:\ntypes · timing · geography", PURPLE),
        ("visualize", "Auto-fires on data answers.\nPlotly dark-theme charts", CORAL),
        ("create-report", "HTML / PNG / CSV reports\npersisted to the sidebar", YELLOW),
        ("311-resolution-rate", "Slash command.\nHistorical resolution %", GREEN),
        ("estimate-complaint", "Slash command.\nSeverity / complexity estimate", CYAN),
    ]
    col_w = Inches(4.0)
    row_h = Inches(2.45)
    gap_x = Inches(0.25)
    gap_y = Inches(0.25)
    start_x = Inches(0.85)
    start_y = Inches(1.8)
    for i, (name, desc, color) in enumerate(cards):
        r, c = divmod(i, 3)
        x = start_x + (col_w + gap_x) * c
        y = start_y + (row_h + gap_y) * r
        add_card(s, x, y, col_w, row_h)
        # top accent bar
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_w, Inches(0.12))
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        add_text(s, name, x + Inches(0.25), y + Inches(0.3), col_w - Inches(0.5), Inches(0.5), size=18, bold=True, color=WHITE, font="Consolas")
        add_text(s, desc, x + Inches(0.25), y + Inches(0.95), col_w - Inches(0.5), Inches(1.3), size=14, color=TEXT)
    return s


def slide_skills_detail(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s, Inches(0.6), PURPLE)
    add_text(s, "Skills in detail", Inches(0.85), Inches(0.55), Inches(12), Inches(0.7), size=36, bold=True, color=WHITE)

    rows = [
        ("download-311-data", "Socrata dataset xwdj-i9he → 311_recent.csv. 100k rows/page. Full download or delta-merge by sr_number.", CYAN),
        ("analyze-311-data", "Top types, departments, status mix, day/hour patterns, ZIPs, districts, resolution times, open backlog, 7-day ASCII bar chart.", PURPLE),
        ("visualize", "Automatic on any data answer with ≥3 points. Plotly dark theme, saved via save_chart, rendered in the artifact panel.", CORAL),
        ("create-report", "Self-contained HTML reports with metric cards + Plotly + narrative. Saved via save_report so they appear in the sidebar.", YELLOW),
        ("311-resolution-rate", "Matches a complaint to a category and returns resolution %, annual volume, reasons for non-resolution, and tips.", GREEN),
        ("estimate-complaint", "Scores severity × complexity × dependencies × resources and outputs a time estimate with confidence + risks.", CYAN),
    ]
    top = Inches(1.45)
    h = Inches(0.85)
    for i, (name, desc, color) in enumerate(rows):
        y = top + Inches(i * 0.92)
        # left color dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.22), Inches(0.35), Inches(0.35))
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        add_text(s, name, Inches(1.35), y + Inches(0.1), Inches(2.9), Inches(0.45), size=16, bold=True, color=WHITE, font="Consolas")
        add_text(s, desc, Inches(4.3), y + Inches(0.12), Inches(8.5), Inches(0.75), size=13, color=TEXT)
    return s


def build():
    prs = new_prs()
    slides = [
        slide_title,
        slide_agenda,
        slide_overview,
        slide_architecture,
        slide_skills_overview,
        slide_skills_detail,
    ]
    for fn in slides:
        fn(prs)
    # footers (skip title)
    total = len(prs.slides)
    for i, s in enumerate(prs.slides, start=1):
        if i == 1:
            continue
        add_footer(s, i, total)
    prs.save(OUT)
    print(f"wrote {OUT} ({total} slides)")


if __name__ == "__main__":
    build()
